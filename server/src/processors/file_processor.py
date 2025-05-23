import os
import logging
from typing import List, Dict, Any
from pathlib import Path
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader
from src.processors.chunking import Chunk, ChunkingStrategy, MarkdownChunker

logger = logging.getLogger(__name__)

class FileProcessor:
    """Processes different file types into chunks."""
    
    def __init__(self, chunking_strategy: ChunkingStrategy = None):
        """Initialize file processor."""
        self.chunking_strategy = chunking_strategy or MarkdownChunker()
        self.supported_extensions = {
            '.md': self._process_markdown,
            '.txt': self._process_text,
            '.docx': self._process_docx,
            '.pptx': self._process_pptx,
            '.pdf': self._process_pdf
        }
    
    def process_file(self, file_path: str) -> List[Chunk]:
        """Process a file into chunks."""
        file_path = Path(file_path)
        
        if not file_path.exists():
            logger.warning(f"File not found: {file_path}")
            return []
            
        extension = file_path.suffix.lower()
        processor = self.supported_extensions.get(extension)
        
        if not processor:
            logger.warning(f"Invalid file type detected for {file_path}")
            return []
            
        try:
            return processor(file_path)
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {str(e)}")
            return []
    
    def _process_markdown(self, file_path: Path) -> List[Chunk]:
        """Process markdown file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
                
            metadata = {
                'source': str(file_path),
                'type': 'markdown'
            }
            
            return self.chunking_strategy.chunk_text(
                content,
                metadata,
                min_chunk_size=100,
                max_chunk_size=512,
                overlap=50
            )
            
        except Exception as e:
            logger.error(f"Error processing markdown file {file_path}: {str(e)}")
            return []
    
    def _process_text(self, file_path: Path) -> List[Chunk]:
        """Process text file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
                
            metadata = {
                'source': str(file_path),
                'type': 'text'
            }
            
            return self.chunking_strategy.chunk_text(
                content,
                metadata,
                min_chunk_size=100,
                max_chunk_size=512,
                overlap=50
            )
            
        except Exception as e:
            logger.error(f"Error processing text file {file_path}: {str(e)}")
            return []
    
    def _process_docx(self, file_path: Path) -> List[Chunk]:
        """Process Word document."""
        try:
            doc = Document(file_path)
            content = "\n".join(p.text for p in doc.paragraphs if p.text)
            
            metadata = {
                'source': str(file_path),
                'type': 'docx'
            }
            
            return self.chunking_strategy.chunk_text(
                content,
                metadata,
                min_chunk_size=100,
                max_chunk_size=512,
                overlap=50
            )
            
        except Exception as e:
            logger.error(f"Error processing Word document {file_path}: {str(e)}")
            return []
    
    def _process_pptx(self, file_path: Path) -> List[Chunk]:
        """Process PowerPoint presentation."""
        try:
            prs = Presentation(file_path)
            content = []
            
            for i, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                if slide_text:
                    content.append(f"Slide {i}:\n" + "\n".join(slide_text))
            
            metadata = {
                'source': str(file_path),
                'type': 'pptx'
            }
            
            return self.chunking_strategy.chunk_text(
                "\n\n".join(content),
                metadata,
                min_chunk_size=100,
                max_chunk_size=512,
                overlap=50
            )
            
        except Exception as e:
            logger.error(f"Error processing PowerPoint file {file_path}: {str(e)}")
            return []
    
    def _process_pdf(self, file_path: Path) -> List[Chunk]:
        """Process PDF document."""
        try:
            reader = PdfReader(file_path)
            content = []
            
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text:
                    content.append(f"Page {i}:\n{text}")
            
            metadata = {
                'source': str(file_path),
                'type': 'pdf'
            }
            
            return self.chunking_strategy.chunk_text(
                "\n\n".join(content),
                metadata,
                min_chunk_size=100,
                max_chunk_size=512,
                overlap=50
            )
            
        except Exception as e:
            logger.error(f"Error processing PDF file {file_path}: {str(e)}")
            return [] 